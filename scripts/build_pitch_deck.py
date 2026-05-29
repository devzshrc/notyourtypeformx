#!/usr/bin/env python3
"""Generate the Schema pitch deck (.pptx)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Palette ──────────────────────────────────────────────────────────────────
BG      = RGBColor(0x0F, 0x17, 0x2A)  # dark navy
PANEL   = RGBColor(0x1E, 0x29, 0x3B)  # panel
ACCENT  = RGBColor(0x7C, 0x3A, 0xED)  # violet
ACCENT2 = RGBColor(0x22, 0xD3, 0xEE)  # cyan
WHITE   = RGBColor(0xF8, 0xFA, 0xFC)
MUTED   = RGBColor(0x94, 0xA3, 0xB8)
GREEN   = RGBColor(0x34, 0xD3, 0x99)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

FONT = "Calibri"


def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(1, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def rect(s, x, y, w, h, color, line=None):
    shp = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    if line: shp.line.color.rgb = line; shp.line.width = Pt(1)
    else: shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=6):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, (txt, size, color, bold) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space); p.space_before = Pt(0)
        r = p.add_run(); r.text = txt
        f = r.font; f.size = Pt(size); f.color.rgb = color; f.bold = bold; f.name = FONT
    return tb


def bullets(s, x, y, w, h, items, size=18, color=WHITE, gap=10):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.alignment = PP_ALIGN.LEFT
        head, body = (it if isinstance(it, tuple) else (it, None))
        r = p.add_run(); r.text = "▸  " + head
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = True; r.font.name = FONT
        if body:
            r2 = p.add_run(); r2.text = "  —  " + body
            r2.font.size = Pt(size - 3); r2.font.color.rgb = MUTED; r2.font.bold = False; r2.font.name = FONT
    return tb


def kicker(s, label):
    rect(s, 0.7, 0.62, 0.12, 0.42, ACCENT)
    text(s, 0.95, 0.55, 8, 0.5, [(label.upper(), 14, ACCENT2, True)])


def heading(s, title):
    text(s, 0.9, 0.95, 11.5, 1.0, [(title, 34, WHITE, True)])


# ── 1. Title ─────────────────────────────────────────────────────────────────
s = slide()
rect(s, 0, 6.95, 13.333, 0.55, ACCENT)
rect(s, 0.9, 2.35, 0.18, 1.6, ACCENT)
text(s, 1.25, 2.25, 11, 1.4, [("Schema", 66, WHITE, True)])
text(s, 1.27, 3.55, 11, 0.8, [("AI-Powered Form Builder", 30, ACCENT2, True)])
text(s, 1.27, 4.45, 10.8, 1.4, [
    ("Build forms by hand or generate them with AI. Collaborate in team", 18, MUTED, False),
    ("workspaces. Publish, embed, collect responses, and track analytics — all type-safe, end to end.", 18, MUTED, False),
])
text(s, 1.27, 6.0, 11, 0.5, [("notyourtypeformx  ·  Pitch Deck", 14, MUTED, True)])

# ── 2. Problem ───────────────────────────────────────────────────────────────
s = slide(); kicker(s, "Problem"); heading(s, "Form tools force a trade-off")
bullets(s, 0.95, 2.05, 11.6, 4.5, [
    ("Polished but locked-in", "Typeform-class UX is beautiful — but expensive, closed, and hard to self-host."),
    ("Powerful but ugly", "Open/free builders are clunky, dated, and break on mobile."),
    ("AI bolted on as an afterthought", "Most tools added a chatbot, not real form generation you can trust."),
    ("No real collaboration", "Teams share one login or copy-paste links; no roles, no workspaces."),
    ("Type drift between front and back end", "API contracts rot; a backend change silently breaks the UI."),
], size=20, gap=16)

# ── 3. Solution ──────────────────────────────────────────────────────────────
s = slide(); kicker(s, "Solution"); heading(s, "One platform, three superpowers")
cards = [
    ("Create", "Hand-build or AI-generate forms with 14 field types, logic jumps & scoring.", ACCENT),
    ("Collaborate", "Team workspaces with OWNER ▸ ADMIN ▸ EDITOR ▸ VIEWER roles & invites.", ACCENT2),
    ("Convert", "Publish, embed, QR — then watch responses & completion rate live.", GREEN),
]
for i, (t, d, c) in enumerate(cards):
    x = 0.95 + i * 4.12
    rect(s, x, 2.3, 3.8, 3.6, PANEL)
    rect(s, x, 2.3, 3.8, 0.14, c)
    text(s, x + 0.3, 2.75, 3.2, 0.7, [(t, 26, c, True)])
    text(s, x + 0.3, 3.6, 3.25, 2.0, [(d, 17, WHITE, False)])

# ── 4. Product / Features ────────────────────────────────────────────────────
s = slide(); kicker(s, "Product"); heading(s, "Everything a form needs — in the box")
col1 = [
    ("14 field types", "text, choice, rating, date, file-style & more"),
    ("Logic jumps & scoring", "branch by answer; build quizzes"),
    ("Drag-drop + bulk import", "reorder, duplicate, paste a list"),
    ("Themes & screens", "welcome / thank-you / custom redirect"),
]
col2 = [
    ("Public slug + iframe embed", "/f/:slug with drop-in embed.js"),
    ("QR codes & draft auto-save", "resume where you left off"),
    ("Access control", "password, response limit, expiry"),
    ("Hidden fields", "capture UTM / query params"),
]
bullets(s, 0.95, 2.1, 6.0, 4.6, col1, size=18, gap=14)
bullets(s, 7.0, 2.1, 5.7, 4.6, col2, size=18, gap=14)

# ── 5. AI ────────────────────────────────────────────────────────────────────
s = slide(); kicker(s, "AI"); heading(s, "AI that builds the form, not a chatbot")
bullets(s, 0.95, 2.05, 7.0, 4.5, [
    ("Prompt → full form", "title, description & 4–12 typed fields in seconds"),
    ("Improve a label", "one-click rewrite to a clearer question"),
    ("Suggest next fields", "context-aware recommendations"),
    ("Prompt-injection guarded", "input sanitized; attacks fall back to a safe default"),
    ("Zod-validated output", "AI never writes unvalidated data to the DB"),
], size=19, gap=15)
rect(s, 8.4, 2.2, 4.2, 3.6, PANEL)
rect(s, 8.4, 2.2, 4.2, 0.14, ACCENT2)
text(s, 8.7, 2.55, 3.7, 0.6, [("Powered by", 14, MUTED, True)])
text(s, 8.7, 3.05, 3.7, 1.6, [
    ("Groq", 30, WHITE, True),
    ("Llama 3.3 70B", 20, ACCENT2, True),
    ("sub-second generation", 15, MUTED, False),
])

# ── 6. Architecture / Moat ───────────────────────────────────────────────────
s = slide(); kicker(s, "How it works"); heading(s, "End-to-end type safety is the moat")
flow = ["Postgres\n+ Drizzle", "Services\n(business logic)", "tRPC router\n(contract)", "Next.js UI\n(typed hooks)"]
for i, label in enumerate(flow):
    x = 0.95 + i * 3.0
    rect(s, x, 2.7, 2.5, 1.5, PANEL, line=ACCENT)
    text(s, x, 2.7, 2.5, 1.5, [(label, 16, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < 3:
        text(s, x + 2.45, 2.7, 0.6, 1.5, [("→", 28, ACCENT2, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
bullets(s, 0.95, 4.7, 11.6, 2.2, [
    ("One contract, zero codegen", "the router defines the API once; the UI infers its types"),
    ("Change backend → UI fails to compile", "no silent drift, fewer production bugs"),
    ("Same router = REST + tRPC + OpenAPI docs", "auto-generated Scalar docs at /docs"),
], size=17, gap=12)

# ── 7. Collaboration ─────────────────────────────────────────────────────────
s = slide(); kicker(s, "Teams"); heading(s, "Built for teams, not just individuals")
bullets(s, 0.95, 2.05, 11.6, 4.5, [
    ("Workspaces with real RBAC", "OWNER ▸ ADMIN ▸ EDITOR ▸ VIEWER, enforced server-side"),
    ("Secure email invitations", "signed tokens, 7-day expiry, email-match verification on accept"),
    ("Full member management", "invite, remove, change role, leave, revoke — with guardrails"),
    ("Forms move freely", "between personal space and shared workspaces"),
    ("Templates & cloning", "publish a form as a template; clone any public form"),
], size=20, gap=16)

# ── 8. Analytics ─────────────────────────────────────────────────────────────
s = slide(); kicker(s, "Analytics"); heading(s, "Know what's working — in real time")
stats = [("Views", ACCENT2), ("Starts", ACCENT), ("Submissions", GREEN), ("Completion %", WHITE)]
for i, (t, c) in enumerate(stats):
    x = 0.95 + i * 3.0
    rect(s, x, 2.5, 2.6, 1.8, PANEL)
    rect(s, x, 2.5, 2.6, 0.12, c)
    text(s, x, 2.85, 2.6, 1.4, [("●", 26, c, True), (t, 17, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
bullets(s, 0.95, 4.8, 11.6, 2.0, [
    ("Live dashboard", "share tab polls every 5 seconds"),
    ("Time-series & admin stats", "trends per form and across all forms"),
    ("Response browser", "date filters, pagination, one-click XLSX export"),
], size=18, gap=12)

# ── 9. Tech stack ────────────────────────────────────────────────────────────
s = slide(); kicker(s, "Stack"); heading(s, "Modern, boring-where-it-counts stack")
rows = [
    ("Monorepo", "Turborepo · pnpm workspaces"),
    ("Frontend", "Next.js 16 · React 19 · Tailwind v4 · Radix · Motion"),
    ("API", "Express 5 · tRPC v11 · TanStack Query · OpenAPI/Scalar"),
    ("Data", "PostgreSQL 15 · Drizzle ORM · Zod v4"),
    ("Auth", "JWT (httpOnly cookies) · bcrypt"),
    ("AI", "Groq SDK · Llama 3.3 70B"),
]
y = 2.1
for i, (k, v) in enumerate(rows):
    c = PANEL if i % 2 == 0 else BG
    rect(s, 0.95, y, 11.4, 0.72, c)
    text(s, 1.2, y, 3.0, 0.72, [(k, 18, ACCENT2, True)], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 4.2, y, 8.0, 0.72, [(v, 17, WHITE, False)], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.78

# ── 10. Differentiation ──────────────────────────────────────────────────────
s = slide(); kicker(s, "Why Schema"); heading(s, "The open, type-safe Typeform alternative")
hdr = ["", "Schema", "Typeform", "OSS builders"]
data = [
    ("Beautiful UX",        "✓", "✓", "—"),
    ("Real AI generation",  "✓", "partial", "—"),
    ("Team RBAC + invites", "✓", "paid tiers", "rare"),
    ("Self-host / own data", "✓", "—", "✓"),
    ("End-to-end type safety", "✓", "n/a", "rare"),
]
cw = [4.2, 2.5, 2.5, 2.5]
x0, y = 0.95, 2.15
# header
cx = x0
for j, h in enumerate(hdr):
    rect(s, cx, y, cw[j], 0.65, ACCENT if j == 1 else PANEL)
    text(s, cx, y, cw[j], 0.65, [(h, 16, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cx += cw[j]
y += 0.7
for row in data:
    cx = x0
    for j, cell in enumerate(row):
        rect(s, cx, y, cw[j], 0.62, BG if j else PANEL)
        col = GREEN if cell == "✓" else (MUTED if cell in ("—", "n/a") else WHITE)
        al = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        pad = cx + (0.25 if j == 0 else 0)
        text(s, pad, y, cw[j] - (0.25 if j == 0 else 0), 0.62,
             [(cell, 15, WHITE if j == 0 else col, j == 0)], align=al, anchor=MSO_ANCHOR.MIDDLE)
        cx += cw[j]
    y += 0.66

# ── 11. Demo / Traction ──────────────────────────────────────────────────────
s = slide(); kicker(s, "Live"); heading(s, "Shipped & running today")
bullets(s, 0.95, 2.05, 7.2, 4.0, [
    ("Full platform live", "auth, forms, AI, workspaces, submissions, analytics"),
    ("Production deploy", "notyourtypeformx.onrender.com"),
    ("Type-checked & linted", "green build across the monorepo"),
    ("Seeded demo data", "forms, responses, workspace, templates"),
], size=19, gap=15)
rect(s, 8.5, 2.2, 4.1, 3.4, PANEL)
rect(s, 8.5, 2.2, 4.1, 0.14, GREEN)
text(s, 8.8, 2.5, 3.6, 0.6, [("Try it — demo login", 15, GREEN, True)])
text(s, 8.8, 3.15, 3.6, 2.0, [
    ("alice@schema.dev", 20, WHITE, True),
    ("Password123!", 18, ACCENT2, True),
    ("", 8, MUTED, False),
    ("/f/customer-feedback", 14, MUTED, False),
    ("/templates", 14, MUTED, False),
])

# ── 12. Roadmap + Ask ────────────────────────────────────────────────────────
s = slide(); kicker(s, "What's next"); heading(s, "Roadmap & the ask")
text(s, 0.95, 2.0, 6.0, 0.5, [("ROADMAP", 16, ACCENT2, True)])
bullets(s, 0.95, 2.5, 6.0, 3.5, [
    ("Payments & paid plans", None),
    ("Webhooks & integrations (Zapier, Slack)", None),
    ("File uploads & e-signatures", None),
    ("Advanced analytics & A/B tests", None),
    ("Custom domains for public forms", None),
], size=17, gap=12)
text(s, 7.2, 2.0, 5.4, 0.5, [("THE ASK", 16, ACCENT2, True)])
rect(s, 7.2, 2.5, 5.4, 3.5, PANEL)
text(s, 7.5, 2.85, 4.9, 3.0, [
    ("Looking for", 16, MUTED, True),
    ("design partners + early teams", 22, WHITE, True),
    ("", 8, MUTED, False),
    ("to shape the roadmap and prove", 16, MUTED, False),
    ("the team-collaboration wedge.", 16, MUTED, False),
])
rect(s, 0, 7.0, 13.333, 0.5, ACCENT)

out = "/Users/dev/Desktop/trpc/Schema-Pitch-Deck.pptx"
prs.save(out)
print("Saved:", out, "| slides:", len(prs.slides._sldIdLst))
